#!/usr/bin/env python3
"""
分析两个PPT的内容逻辑和绘图模式
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import json

def analyze_ppt(pptx_path, name):
    """深度分析单个PPT"""
    prs = Presentation(pptx_path)
    slides_data = []

    for idx, slide in enumerate(prs.slides):
        slide_info = {
            'slide_num': idx + 1,
            'shapes': [],
            'text_content': [],
            'images': [],
            'layout_name': str(slide.slide_layout.name) if slide.slide_layout else 'Unknown',
        }

        for shape in slide.shapes:
            shape_info = {
                'type': shape.shape_type,
                'name': shape.name,
                'position': {
                    'left': shape.left.inches if hasattr(shape, 'left') else 0,
                    'top': shape.top.inches if hasattr(shape, 'top') else 0,
                    'width': shape.width.inches if hasattr(shape, 'width') else 0,
                    'height': shape.height.inches if hasattr(shape, 'height') else 0,
                }
            }

            # 提取文本
            if shape.has_text_frame:
                texts = []
                for para in shape.text_frame.paragraphs:
                    para_text = ''.join([run.text for run in para.runs])
                    if para_text.strip():
                        texts.append({
                            'text': para_text.strip()[:100],
                            'font_size': para.runs[0].font.size.pt if para.runs and para.runs[0].font.size else None,
                            'bold': para.runs[0].font.bold if para.runs else None,
                        })
                if texts:
                    shape_info['texts'] = texts
                    slide_info['text_content'].extend(texts)

            # 提取图片引用
            if shape.shape_type == 13:  # Picture
                shape_info['is_picture'] = True

            # 提取表格
            if shape.shape_type == 19:  # Table
                shape_info['is_table'] = True
                shape_info['table_rows'] = shape.table.rows.__len__()
                shape_info['table_cols'] = len(shape.table.columns)

            slide_info['shapes'].append(shape_info)

        # 统计信息
        slide_info['stats'] = {
            'total_shapes': len(slide.shapes),
            'text_boxes': sum(1 for s in slide.shapes if s.has_text_frame),
            'images': sum(1 for s in slide.shapes if s.shape_type == 13),
            'tables': sum(1 for s in slide.shapes if s.shape_type == 19),
            'lines': sum(1 for s in slide.shapes if s.shape_type == 1),
            'rectangles': sum(1 for s in slide.shapes if s.shape_type == 9),
        }

        slides_data.append(slide_info)

    return slides_data

def extract_content_patterns(slides_data, name):
    """提取内容模式"""
    patterns = {
        'name': name,
        'total_slides': len(slides_data),
        'title_slides': [],
        'content_slides': [],
        'common_elements': [],
        'page_structures': [],
    }

    # 分析每页结构
    for slide in slides_data:
        stats = slide['stats']
        texts = slide['text_content']

        # 识别标题页 vs 内容页
        if stats['text_boxes'] <= 5 and any('标题' in t['text'] or '封面' in t['text'] or '感谢' in t['text'] for t in texts if isinstance(t, dict)):
            patterns['title_slides'].append(slide['slide_num'])
        elif stats['text_boxes'] > 2:
            patterns['content_slides'].append(slide['slide_num'])

        # 记录页面结构
        page_structure = {
            'slide_num': slide['slide_num'],
            'shape_count': stats['total_shapes'],
            'text_boxes': stats['text_boxes'],
            'images': stats['images'],
            'tables': stats['tables'],
            'layout': slide['layout_name'],
        }
        patterns['page_structures'].append(page_structure)

    # 统计常见元素
    all_shapes = []
    for slide in slides_data:
        all_shapes.extend([s['type'] for s in slide['shapes']])

    patterns['shape_distribution'] = {str(k): v for k, v in sorted(__import__('collections').Counter(all_shapes).items(), key=lambda x: -x[1])[:5]}

    return patterns

def extract_writing_logic(slides_data, name):
    """提取写作逻辑模式"""
    writing_patterns = {
        'name': name,
        'content_types': [],
        'structure_patterns': [],
        'narrative_flows': [],
    }

    for slide in slides_data:
        texts = slide['text_content']
        if not texts:
            continue

        # 提取首行作为标题参考
        main_titles = [t['text'] for t in texts if isinstance(t, dict) and t.get('font_size', 0) and t['font_size'] > 18]
        all_texts = [t['text'] for t in texts if isinstance(t, dict)]

        if main_titles:
            writing_patterns['content_types'].append({
                'slide': slide['slide_num'],
                'main_title': main_titles[0][:50] if main_titles else '',
                'bullet_count': len([t for t in all_texts if len(t) < 50]),
                'paragraph_count': len([t for t in all_texts if len(t) > 50]),
            })

    return writing_patterns

def main():
    pptx_files = [
        ('/Users/ciondlin/Downloads/长亭安服主打胶片- v2.2-0427.pptx', '长亭安服主打胶片'),
        ('/Users/ciondlin/Downloads/2025HW总结-25.09.004.pptx', '2025HW总结'),
    ]

    all_analysis = {}

    for pptx_path, name in pptx_files:
        if not os.path.exists(pptx_path):
            print(f"文件不存在: {pptx_path}")
            continue

        print(f"\n{'='*60}")
        print(f"分析: {name}")
        print(f"{'='*60}")

        slides_data = analyze_ppt(pptx_path, name)
        patterns = extract_content_patterns(slides_data, name)
        writing = extract_writing_logic(slides_data, name)

        print(f"\n总页数: {patterns['total_slides']}")
        print(f"标题页: {patterns['title_slides']}")
        print(f"内容页: {patterns['content_slides'][:10]}..." if len(patterns['content_slides']) > 10 else f"内容页: {patterns['content_slides']}")
        print(f"形状分布: {patterns['shape_distribution']}")

        # 打印前5页的内容摘要
        print(f"\n前5页内容摘要:")
        for slide in slides_data[:5]:
            texts = slide['text_content']
            if texts:
                print(f"  第{slide['slide_num']}页: {[t['text'][:30] for t in texts[:3] if isinstance(t, dict)]}")

        all_analysis[name] = {
            'patterns': patterns,
            'writing': writing,
        }

    # 保存分析结果
    output_path = '/Users/ciondlin/skills/ppt-master/templates/layouts/security_service/analysis_report.md'
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('# PPT编辑逻辑与绘图模式分析报告\n\n')
        for name, data in all_analysis.items():
            f.write(f'\n## {name}\n\n')
            f.write(f'### 页面结构统计\n')
            f.write(f'- 总页数: {data["patterns"]["total_slides"]}\n')
            f.write(f'- 标题页: {data["patterns"]["title_slides"]}\n')
            f.write(f'- 内容页数: {len(data["patterns"]["content_slides"])}\n\n')

            f.write(f'### 形状分布\n')
            for shape_type, count in data['patterns']['shape_distribution'].items():
                f.write(f'- {shape_type}: {count}\n')
            f.write('\n')

            f.write(f'### 页面结构详情\n')
            for ps in data['patterns']['page_structures'][:10]:
                f.write(f'- 第{ps["slide_num"]}页: 形状{ps["shape_count"]}个, 文本框{ps["text_boxes"]}个, 图片{ps["images"]}个, 表格{ps["tables"]}个, 布局:{ps["layout"]}\n')
            f.write('\n')

            f.write(f'### 内容类型示例\n')
            for ct in data['writing']['content_types'][:8]:
                f.write(f'- 第{ct["slide"]}页「{ct["main_title"]}」: {ct["bullet_count"]}个要点, {ct["paragraph_count"]}段文字\n')

    print(f"\n\n分析报告已保存到: {output_path}")

if __name__ == '__main__':
    main()
