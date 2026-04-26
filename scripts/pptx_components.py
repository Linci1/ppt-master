#!/usr/bin/env python3
"""
pptx_components.py
==================
python-pptx 直接生成 body 页的组件库。
不再以 SVG 为中间格式，直接控制每个 shape 的属性。

组件分两层：
  原子组件  — add_card / add_rect / add_textbox / add_brand_header / add_brand_footer
  页型函数  — build_card_grid_page / build_timeline_page / build_data_table_page
               build_two_column_page / build_text_section_page

使用流程：
  1. 用 reference_analyzer.py 从参考稿提取 layout_index.json
  2. 用 pptx_native_executor.py 生成页型 Python 代码
  3. 执行代码 → PPTX
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path
from typing import List, Dict, Tuple, Optional

# ─────────────────────────────────────────────────────────────────
# Theme 配置（可按模板覆盖）
# ─────────────────────────────────────────────────────────────────

class Theme:
    """主题配置。创建 Presentation 前设置 Theme.* 类变量。"""

    CANVAS_W = Inches(13.333)
    CANVAS_H = Inches(7.5)

    # 字号层级
    FS_HERO    = 28   # 页面大标题
    FS_SUB     = 13   # 副标题/描述
    FS_CARD_T  = 16   # 卡片标题
    FS_BODY    = 13   # 正文
    FS_SMALL   = 12   # 页码/注释
    FS_TABLE_H = 14   # 表头
    FS_TABLE_C = 12   # 表格内容

    # 边距
    MARGIN_LEFT   = 0.42   # 内容区左边界 (in)
    MARGIN_RIGHT  = 0.42
    MARGIN_TOP    = 1.20   # 内容区上边 y（品牌标题区以下）
    MARGIN_BOTTOM = 6.90   # 内容区下边 y

    # 配色盘（子类覆盖）
    PRIMARY   = RGBColor(0x7B, 0xBD, 0x4A)   # 品牌绿
    BG        = RGBColor(0xFF, 0xFF, 0xFF)   # 页面背景
    TEXT_HERO = RGBColor(0x1A, 0x1A, 0x1A)   # 深色主文字
    TEXT_SEC  = RGBColor(0x66, 0x66, 0x66)   # 灰色副文字
    TEXT_LITE = RGBColor(0xA6, 0xA6, 0xA6)   # 浅灰（页码）
    DIVIDER   = RGBColor(0xE0, 0xE0, 0xE0)   # 分隔线
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
    CARD_DEEP  = RGBColor(0x3C, 0x74, 0x71)  # 深青色（默认深卡背景）

    @classmethod
    def card_dark(cls):
        """深色卡片风格（深青背景 + 品牌绿标题条）"""
        return {
            'bg':        cls.CARD_DEEP,
            'bar':       cls.PRIMARY,
            'title':     cls.WHITE,
            'body':      RGBColor(0xF5, 0xF5, 0xF5),
        }

    @classmethod
    def card_light(cls):
        """浅色卡片风格（浅灰背景 + 绿色标题条）"""
        return {
            'bg':        RGBColor(0xF5, 0xF7, 0xFA),
            'bar':       cls.PRIMARY,
            'title':     cls.WHITE,
            'body':      cls.TEXT_HERO,
        }


# 快捷配色（chaitin_anfu 默认）
C = Theme()

# ─────────────────────────────────────────────────────────────────
# 原子组件
# ─────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb: RGBColor, line_rgb: RGBColor = None):
    """添加矩形 shape"""
    shape = slide.shapes.add_shape(1, int(x), int(y), int(w), int(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_run(paragraph, text: str, size_pt: float,
                  color_rgb: RGBColor, bold: bool = False,
                  font_name: str = 'Microsoft YaHei') -> None:
    """向段落添加一个 run"""
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color_rgb
    run.font.name = font_name


def add_textbox(slide, x, y, w, h,
                lines: List[Tuple[str, float, RGBColor, bool]],
                align=PP_ALIGN.LEFT,
                wrap=True) -> None:
    """
    添加 TextBox，支持多行。

    lines: List[(text, size_pt, color_rgb, bold)]
    """
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, (text, size_pt, color_rgb, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(2)
        add_text_run(p, text, size_pt, color_rgb, bold)


def add_card(slide,
             x: float, y: float,
             card_w: float, card_h: float,
             title: str,
             body_lines: List[str],
             style: Dict = None) -> None:
    """
    添加一张卡片（标题条 + 内容区）。

    style: dict with keys bg/bar/title/body (RGBColor)
           or 'dark'/'light' string (use Theme.card_dark/card_light)
    """
    x, y, card_w, card_h = Inches(x), Inches(y), Inches(card_w), Inches(card_h)

    if isinstance(style, str):
        style = Theme.card_dark() if style == 'dark' else Theme.card_light()

    bg_c     = style['bg']
    bar_c    = style['bar']
    title_c  = style['title']
    body_c   = style['body']
    bar_h    = Inches(0.52)
    pad_x    = Inches(0.16)
    pad_top  = Inches(0.10)

    # 1. 卡片背景
    add_rect(slide, x, y, card_w, card_h, bg_c)
    # 2. 标题条
    add_rect(slide, x, y, card_w, bar_h, bar_c)
    # 3. 标题文字
    tb_t = slide.shapes.add_textbox(
        x + pad_x, y + pad_top,
        card_w - pad_x * 2, bar_h - pad_top
    )
    tf = tb_t.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(0); p.space_after = Pt(0)
    add_text_run(p, title, C.FS_CARD_T, title_c, bold=True)

    # 4. 内容文字
    content_top = y + bar_h + Inches(0.22)
    content_h   = card_h - bar_h - Inches(0.22)
    tb_b = slide.shapes.add_textbox(
        x + pad_x, content_top,
        card_w - pad_x * 2, content_h
    )
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf_b.paragraphs[0] if i == 0 else tf_b.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0)
        p.space_after  = Pt(6)
        add_text_run(p, line, C.FS_BODY, body_c, bold=False)


def add_brand_header(slide,
                     page_title: str,
                     page_num: str = '4',
                     logo_path: str = None) -> None:
    """
    品牌锁定 Header：
      顶部 0.04" 绿色条 + 页面标题 + 分隔线 + 页码 + logo（可选）
    """
    # 顶部绿条
    add_rect(slide, 0, 0, C.CANVAS_W, Inches(0.04), C.PRIMARY)
    # 页面大标题
    tb = slide.shapes.add_textbox(
        int(Inches(0.34)), int(Inches(0.30)),
        int(Inches(5)), int(Inches(0.45))
    )
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(0); p.space_after = Pt(0)
    add_text_run(p, page_title, C.FS_HERO, C.PRIMARY, bold=True)
    # 分隔线
    add_rect(slide,
             int(Inches(0.34)), int(Inches(0.78)),
             int(Inches(12.5)), int(Inches(0.02)),
             C.DIVIDER)
    # 页码
    tb_pn = slide.shapes.add_textbox(
        int(Inches(12.2)), int(Inches(0.38)),
        int(Inches(0.8)), int(Inches(0.3))
    )
    p_pn = tb_pn.text_frame.paragraphs[0]
    p_pn.alignment = PP_ALIGN.RIGHT
    p_pn.space_before = Pt(0); p_pn.space_after = Pt(0)
    add_text_run(p_pn, page_num, C.FS_SMALL, C.TEXT_LITE, bold=False)
    # Logo
    if logo_path and Path(logo_path).exists():
        slide.shapes.add_picture(
            logo_path,
            int(Inches(11.45)), int(Inches(0.25)),
            int(Inches(1.1)), int(Inches(0.32))
        )


def add_brand_footer(slide) -> None:
    """品牌锁定 Footer：底部绿色条"""
    add_rect(slide,
             0, int(Inches(7.46)),
             C.CANVAS_W, int(Inches(0.04)),
             C.PRIMARY)


def add_timeline_item(slide,
                      x: float, y: float,
                      w: float, h: float,
                      time_label: str,
                      title: str,
                      description: str,
                      dot_color: RGBColor = None,
                      line_color: RGBColor = None) -> None:
    """
    添加一个时间线节点。

    layout:
      [dot]───[time_label]    y
              [title]        bold
              [description]  normal
    """
    x, y, w, h = Inches(x), Inches(y), Inches(w), Inches(h)
    dot_color = dot_color or C.PRIMARY
    line_color = line_color or C.DIVIDER

    dot_size = Inches(0.20)
    dot_x = x
    dot_y = y + Inches(0.10)

    # 垂直连接线（向上延伸）
    if line_color:
        line = slide.shapes.add_shape(
            1, dot_x + dot_size // 2 - int(Pt(1)),
            dot_y - Inches(0.60),
            int(Pt(2)), Inches(0.60)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = line_color
        line.line.fill.background()

    # 圆点
    add_rect(slide, dot_x, dot_y, dot_size, dot_size,
             dot_color)

    # 时间标签（右侧）
    tb_time = slide.shapes.add_textbox(
        x + dot_size + Inches(0.12),
        y + Inches(0.04),
        Inches(1.4), Inches(0.28)
    )
    p = tb_time.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(0); p.space_after = Pt(0)
    add_text_run(p, time_label, C.FS_SMALL, C.TEXT_LITE, bold=False)

    # 标题
    tb_title = slide.shapes.add_textbox(
        x + dot_size + Inches(0.12),
        y + Inches(0.30),
        Inches(3.5), Inches(0.30)
    )
    p = tb_title.text_frame.paragraphs[0]
    p.space_before = Pt(0); p.space_after = Pt(0)
    add_text_run(p, title, C.FS_BODY, C.TEXT_HERO, bold=True)

    # 描述
    tb_desc = slide.shapes.add_textbox(
        x + dot_size + Inches(0.12),
        y + Inches(0.62),
        Inches(4.5), Inches(0.40)
    )
    tf = tb_desc.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.space_before = Pt(0); p.space_after = Pt(0)
    add_text_run(p, description, C.FS_SMALL, C.TEXT_SEC, bold=False)


def add_table(slide,
              x: float, y: float,
              w: float,
              headers: List[str],
              rows: List[List[str]],
              col_widths: List[float] = None) -> None:
    """
    添加带样式的表格。

    headers: 表头文字列表
    rows: 数据行列表
    col_widths: 每列宽度（in），默认均分
    """
    x, y, w = Inches(x), Inches(y), Inches(w)
    n_cols = len(headers)
    col_widths = col_widths or [w / n_cols] * n_cols
    row_h = Inches(0.42)
    header_h = Inches(0.48)

    # 计算总高度
    total_h = header_h + row_h * len(rows)

    # 外框
    add_rect(slide, x, y, w, total_h,
             RGBColor(0xE0, 0xE0, 0xE0))

    # 表头背景
    add_rect(slide, x, y, w, header_h, C.PRIMARY)

    # 绘制表头
    cx = x
    for hi, (hdr, cw) in enumerate(zip(headers, col_widths)):
        cw = Inches(cw)
        tb = slide.shapes.add_textbox(
            cx + Inches(0.10), y + Inches(0.10),
            cw - Inches(0.20), header_h - Inches(0.10)
        )
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0); p.space_after = Pt(0)
        add_text_run(p, hdr, C.FS_TABLE_H, C.WHITE, bold=True)

        # 表头竖线
        if hi > 0:
            add_rect(slide,
                     cx, y,
                     int(Pt(0.5)), header_h,
                     RGBColor(0x5A, 0x9E, 0x3A))  # 深绿分隔线

        cx += cw

    # 绘制数据行
    for ri, row in enumerate(rows):
        ry = y + header_h + row_h * ri
        bg = C.WHITE if ri % 2 == 0 else RGBColor(0xF5, 0xF7, 0xFA)

        # 行背景
        add_rect(slide, x, ry, w, row_h, bg)

        # 单元格
        cx = x
        for ci, (cell, cw) in enumerate(zip(row, col_widths)):
            cw = Inches(cw)
            tb = slide.shapes.add_textbox(
                cx + Inches(0.10), ry + Inches(0.08),
                cw - Inches(0.20), row_h - Inches(0.08)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(0); p.space_after = Pt(0)
            add_text_run(p, str(cell), C.FS_TABLE_C,
                         C.TEXT_HERO, bold=False)

            # 单元格竖线
            if ci > 0:
                add_rect(slide,
                         cx, ry,
                         int(Pt(0.5)), row_h,
                         RGBColor(0xE0, 0xE0, 0xE0))
            cx += cw

        # 行横线
        add_rect(slide,
                 x, ry + row_h - int(Pt(0.5)),
                 w, int(Pt(0.5)),
                 C.DIVIDER)


# ─────────────────────────────────────────────────────────────────
# 页型组装函数
# ─────────────────────────────────────────────────────────────────

def build_card_grid_page(prs: Presentation,
                         title: str,
                         subtitle: str = '',
                         cards: List[Dict] = None,
                         page_num: str = '4',
                         logo_path: str = None,
                         card_style = 'dark') -> None:
    """
    2×2 卡片网格页面。

    cards: [
        {'title': str, 'lines': [str, ...]},
        ...
    ]  共 4 项

    card_style: 'dark' / 'light' / dict (Theme.card_dark() 等)
    """
    cards = cards or []
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    add_brand_header(slide, title, page_num, logo_path)

    # 副标题
    if subtitle:
        tb_sub = slide.shapes.add_textbox(
            int(Inches(0.34)), int(Inches(0.82)),
            int(Inches(8)), int(Inches(0.25))
        )
        p = tb_sub.text_frame.paragraphs[0]
        p.space_before = Pt(0); p.space_after = Pt(0)
        add_text_run(p, subtitle, C.FS_SUB, C.TEXT_SEC, bold=False)

    # 卡片网格
    gx0 = 0.42;  gy0 = 1.25
    cw  = 5.79;  ch  = 2.40
    gx  = 0.50;  gy  = 0.45

    positions = [
        (gx0,           gy0),
        (gx0 + cw + gx, gy0),
        (gx0,           gy0 + ch + gy),
        (gx0 + cw + gx, gy0 + ch + gy),
    ]

    style = (Theme.card_dark() if card_style == 'dark'
              else Theme.card_light() if card_style == 'light'
              else card_style)

    for card, (x, y) in zip(cards, positions):
        add_card(slide, x, y, cw, ch,
                 card['title'], card.get('lines', []),
                 style=style)

    add_brand_footer(slide)


def build_timeline_page(prs: Presentation,
                        title: str,
                        subtitle: str = '',
                        items: List[Dict] = None,
                        page_num: str = '4',
                        logo_path: str = None,
                        accent_color: RGBColor = None) -> None:
    """
    时间线页面。

    items: [
        {'time': str, 'title': str, 'description': str},
        ...
    ]
    """
    items = items or []
    accent_color = accent_color or C.PRIMARY
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_brand_header(slide, title, page_num, logo_path)

    if subtitle:
        tb_sub = slide.shapes.add_textbox(
            int(Inches(0.34)), int(Inches(0.82)),
            int(Inches(8)), int(Inches(0.25))
        )
        p = tb_sub.text_frame.paragraphs[0]
        p.space_before = Pt(0); p.space_after = Pt(0)
        add_text_run(p, subtitle, C.FS_SUB, C.TEXT_SEC, bold=False)

    # 内容区起始 y
    start_y = 1.25
    row_h   = 1.10   # 每行高度
    card_w  = 5.50
    gap     = 0.60   # 列间距

    for i, item in enumerate(items):
        ri = i % 3   # 行内索引
        ci = i // 3  # 列索引
        x  = 0.42 + ci * (card_w + gap)
        y  = start_y + ri * row_h

        add_timeline_item(
            slide, x, y, card_w, row_h,
            time_label=item.get('time', ''),
            title=item.get('title', ''),
            description=item.get('description', ''),
            dot_color=accent_color,
        )

    add_brand_footer(slide)


def build_data_table_page(prs: Presentation,
                           title: str,
                           subtitle: str = '',
                           headers: List[str] = None,
                           rows: List[List[str]] = None,
                           page_num: str = '4',
                           logo_path: str = None,
                           col_widths: List[float] = None) -> None:
    """
    数据表格页面。

    headers: 表头列表
    rows: 数据行列表
    col_widths: 每列宽度（in）
    """
    headers = headers or []
    rows    = rows or []
    slide   = prs.slides.add_slide(prs.slide_layouts[6])

    add_brand_header(slide, title, page_num, logo_path)

    if subtitle:
        tb_sub = slide.shapes.add_textbox(
            int(Inches(0.34)), int(Inches(0.82)),
            int(Inches(8)), int(Inches(0.25))
        )
        p = tb_sub.text_frame.paragraphs[0]
        p.space_before = Pt(0); p.space_after = Pt(0)
        add_text_run(p, subtitle, C.FS_SUB, C.TEXT_SEC, bold=False)

    # 表格位置
    tbl_x = 0.42
    tbl_y = 1.30
    tbl_w = 12.50

    add_table(slide, tbl_x, tbl_y, tbl_w,
              headers, rows, col_widths)

    add_brand_footer(slide)


def build_two_column_page(prs: Presentation,
                          title: str,
                          subtitle: str = '',
                          left: Dict = None,
                          right: Dict = None,
                          page_num: str = '4',
                          logo_path: str = None) -> None:
    """
    双栏对比页面。

    left/right: {
        'heading': str,   # 栏目标题
        'bullets': [str, ...],  # 要点列表
    }
    """
    left  = left or {}
    right = right or {}
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_brand_header(slide, title, page_num, logo_path)

    if subtitle:
        tb_sub = slide.shapes.add_textbox(
            int(Inches(0.34)), int(Inches(0.82)),
            int(Inches(8)), int(Inches(0.25))
        )
        p = tb_sub.text_frame.paragraphs[0]
        p.space_before = Pt(0); p.space_after = Pt(0)
        add_text_run(p, subtitle, C.FS_SUB, C.TEXT_SEC, bold=False)

    # 内容区
    col_w  = 5.90
    col_h  = 4.80
    start_x = 0.42
    start_y = 1.25
    gap     = 0.58

    for ci, col in enumerate([left, right]):
        x = start_x + ci * (col_w + gap)
        bg_c = Theme.card_dark()['bg'] if ci == 0 else Theme.card_light()['bg']

        # 栏背景
        add_rect(slide, int(Inches(x)), int(Inches(start_y)),
                 int(Inches(col_w)), int(Inches(col_h)), bg_c)

        # 顶部标题条
        bar_h = Inches(0.50)
        add_rect(slide, int(Inches(x)), int(Inches(start_y)),
                 int(Inches(col_w)), bar_h, C.PRIMARY)

        # 栏目标题
        tb_h = slide.shapes.add_textbox(
            int(Inches(x + 0.18)), int(Inches(start_y + 0.08)),
            int(Inches(col_w - 0.36)), bar_h - Inches(0.08)
        )
        tf = tb_h.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0); p.space_after = Pt(0)
        add_text_run(p, col.get('heading', ''),
                     C.FS_CARD_T, C.WHITE, bold=True)

        # 要点列表
        bullets = col.get('bullets', [])
        bullet_y = start_y + 0.65
        for bi, bullet in enumerate(bullets):
            # 绿色圆点
            dot = slide.shapes.add_shape(
                1,
                int(Inches(x + 0.18)),
                int(Inches(bullet_y + 0.12)),
                int(Inches(0.12)), int(Inches(0.12))
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = C.PRIMARY
            dot.line.fill.background()

            tb_b = slide.shapes.add_textbox(
                int(Inches(x + 0.40)),
                int(Inches(bullet_y)),
                int(Inches(col_w - 0.60)),
                int(Inches(0.50))
            )
            tf_b = tb_b.text_frame
            tf_b.word_wrap = True
            p = tf_b.paragraphs[0]
            p.space_before = Pt(0); p.space_after = Pt(0)
            add_text_run(p, bullet, C.FS_BODY, C.TEXT_HERO, bold=False)
            bullet_y += 0.70

    add_brand_footer(slide)


def build_text_section_page(prs: Presentation,
                            title: str,
                            subtitle: str = '',
                            sections: List[Dict] = None,
                            page_num: str = '4',
                            logo_path: str = None) -> None:
    """
    纯文本段落页面（适用于政策说明、术语解释等）。

    sections: [
        {'heading': str, 'body': str},   # 小标题 + 正文段落
        ...
    ]
    """
    sections = sections or []
    slide    = prs.slides.add_slide(prs.slide_layouts[6])

    add_brand_header(slide, title, page_num, logo_path)

    if subtitle:
        tb_sub = slide.shapes.add_textbox(
            int(Inches(0.34)), int(Inches(0.82)),
            int(Inches(8)), int(Inches(0.25))
        )
        p = tb_sub.text_frame.paragraphs[0]
        p.space_before = Pt(0); p.space_after = Pt(0)
        add_text_run(p, subtitle, C.FS_SUB, C.TEXT_SEC, bold=False)

    # 内容
    cur_y = 1.30
    sec_gap = 0.55   # 段落间距

    for sec in sections:
        heading = sec.get('heading', '')
        body    = sec.get('body', '')

        # 小标题
        if heading:
            tb_h = slide.shapes.add_textbox(
                int(Inches(0.42)), int(Inches(cur_y)),
                int(Inches(12)), int(Inches(0.35))
            )
            p = tb_h.text_frame.paragraphs[0]
            p.space_before = Pt(0); p.space_after = Pt(0)
            add_text_run(p, heading, C.FS_BODY, C.PRIMARY, bold=True)
            cur_y += 0.38

        # 正文（自动换行）
        tb_b = slide.shapes.add_textbox(
            int(Inches(0.42)), int(Inches(cur_y)),
            int(Inches(12)), int(Inches(1.20))
        )
        tf = tb_b.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.space_before = Pt(0); p.space_after = Pt(4)
        add_text_run(p, body, C.FS_BODY, C.TEXT_HERO, bold=False)
        cur_y += sec_gap

    add_brand_footer(slide)


# ─────────────────────────────────────────────────────────────────
# 快捷入口：从 JSON 描述生成（供 executor 调用）
# ─────────────────────────────────────────────────────────────────

def page_type_from_string(type_str: str):
    """把字符串页型名映射到函数"""
    MAP = {
        'card_grid':      build_card_grid_page,
        'timeline':       build_timeline_page,
        'data_table':    build_data_table_page,
        'two_column':     build_two_column_page,
        'text_section':   build_text_section_page,
    }
    return MAP.get(type_str, build_card_grid_page)


# ─────────────────────────────────────────────────────────────────
# CLI 验证
# ─────────────────────────────────────────────────────────────────
if __name__ == '__main__':
    from pptx import Presentation as P

    prs = P()
    prs.slide_width  = C.CANVAS_W
    prs.slide_height = C.CANVAS_H

    # 1. 卡片网格页
    build_card_grid_page(prs,
        title='攻击路径总览',
        subtitle='Attack Path Overview',
        cards=[
            {'title': '边界突破', 'lines': ['ThinkPHP RCE', 'Log4j2 JNDI']},
            {'title': '内网横移', 'lines': ['社工钓鱼', '横向扩展']},
            {'title': '社工钓鱼', 'lines': ['微信投递简历', '木马触发']},
            {'title': '核心资产', 'lines': ['AWS控制权', '域控/MySQL']},
        ],
        page_num='4',
        card_style='dark')

    # 2. 时间线页
    build_timeline_page(prs,
        title='攻击时间线',
        subtitle='Attack Timeline | 2025-04-21 ~ 2025-04-28',
        items=[
            {'time': 'Day 1', 'title': '边界发现', 'description': 'Nmap扫描发现VPN入口'},
            {'time': 'Day 2', 'title': '漏洞利用', 'description': 'ThinkPHP RCE获取首台服务器'},
            {'time': 'Day 3', 'title': '权限提升', 'description': 'sudo提权获取root'},
            {'time': 'Day 5', 'title': '内网横移', 'description': '利用hash传递横向'},
            {'time': 'Day 7', 'title': '目标达成', 'description': '域控沦陷'},
        ],
        page_num='5')

    # 3. 数据表格页
    build_data_table_page(prs,
        title='漏洞清单',
        subtitle='Vulnerability Assessment',
        headers=['漏洞编号', '严重程度', '影响组件', '修复建议'],
        rows=[
            ['CVE-2024-001', '高危', 'ThinkPHP 5.x', '升级至5.1.41+'],
            ['CVE-2024-002', '严重', 'Apache Log4j2', '升级至2.17.0'],
            ['CVE-2024-003', '中危', 'SSH服务', '禁用弱口令策略'],
        ],
        col_widths=[2.0, 1.8, 2.5, 3.2],
        page_num='6')

    # 4. 双栏对比页
    build_two_column_page(prs,
        title='攻防对比',
        subtitle='Offensive vs Defensive Analysis',
        left={'heading': '攻击方', 'bullets': ['0day漏洞利用', '社工钓鱼', '内网横移']},
        right={'heading': '防守方', 'bullets': ['边界防护不足', '员工安全意识弱', '日志审计缺失']},
        page_num='7')

    # 5. 文本段落页
    build_text_section_page(prs,
        title='术语说明',
        subtitle='Glossary of Terms',
        sections=[
            {'heading': 'Red Team（红队）',
             'body': '模拟真实攻击者行为的专业团队，通过各种攻击手段验证目标网络的防御能力。'},
            {'heading': 'Purple Team（紫队）',
             'body': '红蓝双方的协作模式，通过攻击复盘和防御改进提升整体安全水平。'},
        ],
        page_num='8')

    out = '/tmp/pptx_components_demo.pptx'
    prs.save(out)
    print(f'✓ 5种页型演示: {out}')
    print(f'  Slide 1: card_grid | Slide 2: timeline | Slide 3: data_table')
    print(f'  Slide 4: two_column | Slide 5: text_section')
