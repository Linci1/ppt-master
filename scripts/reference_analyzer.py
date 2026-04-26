#!/usr/bin/env python3
"""
reference_analyzer.py
=====================
从参考 PPTX 文件提取结构化布局数据，输出 layout_index.json。

使用方式：
  python3 reference_analyzer.py <pptx_file> [--output <json_path>]

输出：
  <pptx_file>_layout_index.json
  ├── metadata       # 画布尺寸、页数、主题色
  ├── pages[]        # 每页元素（矩形/文字/表格）的完整结构
  ├── brand_layout   # 品牌锁定区（header/footer）
  ├── content_layouts[]  # 每页内容区的布局类型推断
  └── component_hints[]  # 建议的页型函数

依赖：python-pptx, lxml
"""

import json, sys, zipfile, argparse, textwrap
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from typing import Dict, List, Any, Optional

I = 914400  # 1 inch = 914400 EMU

# ─────────────────────────────────────────────────────────────────
# XML 解析
# ─────────────────────────────────────────────────────────────────

NS = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

def parse_pptx(pptx_path: str) -> Dict:
    """返回 {'metadata', 'slides', 'theme_colors'}"""
    prs = Presentation(pptx_path)

    # 画布尺寸
    cw = prs.slide_width
    ch = prs.slide_height
    canvas_w_in = cw / I
    canvas_h_in = ch / I

    # 幻灯片数量
    n_slides = len(prs.slides)

    # 尝试提取主题色（从第1页第一个形状）
    theme_colors = extract_theme_colors(pptx_path)

    slides_data = []
    with zipfile.ZipFile(pptx_path) as z:
        for i, slide in enumerate(prs.slides, 1):
            slide_file = f'ppt/slides/slide{i}.xml'
            try:
                xml_bytes = z.read(slide_file)
            except KeyError:
                slides_data.append({'index': i, 'error': 'slide file not found'})
                continue

            elements = extract_elements(xml_bytes)
            slides_data.append({
                'index':    i,
                'file':     slide_file,
                'elements': elements,
            })

    return {
        'metadata': {
            'source_pptx': str(pptx_path),
            'canvas_w_in': round(canvas_w_in, 3),
            'canvas_h_in': round(canvas_h_in, 3),
            'n_slides': n_slides,
            'theme_colors': theme_colors,
        },
        'slides': slides_data,
    }


def extract_theme_colors(pptx_path: str) -> Dict[str, str]:
    """从幻灯片母版/主题提取主色"""
    colors = {}
    try:
        with zipfile.ZipFile(pptx_path) as z:
            # 读主题
            theme_xml = z.read('ppt/theme/theme1.xml')
            root = etree.fromstring(theme_xml)
            # 读 clrScheme 的第一个颜色（通常是dk1或lt1旁的自定义色）
            scheme = root.find('.//a:clrScheme', NS)
            if scheme is not None:
                for child in scheme:
                    tag = child.tag.split('}')[-1]
                    srgb = child.find('a:srgbClr', NS)
                    if srgb is not None:
                        colors[tag] = srgb.get('val')
    except Exception:
        pass
    return colors


def extract_elements(xml_bytes: bytes) -> List[Dict]:
    """从 slide XML 提取所有 shape 的坐标/样式/文字"""
    root = etree.fromstring(xml_bytes)
    shapes = root.findall('.//p:sp', NS)
    pic_shapes = root.findall('.//p:pic', NS)
    grp_shapes = root.findall('.//p:grpSp', NS)
    all_shapes = shapes + pic_shapes + grp_shapes

    elements = []
    for shape in all_shapes:
        el = parse_shape(shape)
        if el:
            elements.append(el)

    # 按 y 坐标排序（从上到下）
    elements.sort(key=lambda e: (e['y_in'], e['x_in']))
    return elements


def parse_shape(shape) -> Optional[Dict]:
    """解析单个 shape，返回结构化数据或 None"""
    xfrm = shape.find('.//a:xfrm', NS)
    if xfrm is None:
        return None

    off = xfrm.find('a:off', NS)
    ext = xfrm.find('a:ext', NS)
    if off is None or ext is None:
        return None

    x_emu = int(off.get('x', 0))
    y_emu = int(off.get('y', 0))
    cx_emu = int(ext.get('cx', 0))
    cy_emu = int(ext.get('cy', 0))

    x_in = round(x_emu / I, 4)
    y_in = round(y_emu / I, 4)
    w_in = round(cx_emu / I, 4)
    h_in = round(cy_emu / I, 4)

    # 填充色
    solidFill = shape.find('.//a:solidFill/a:srgbClr', NS)
    fill_color = solidFill.get('val', '').upper() if solidFill is not None else ''

    # gradFill（渐变）
    gradFill = shape.find('.//a:gradFill', NS)
    grad_info = None
    if gradFill is not None:
        clr = gradFill.find('.//a:srgbClr', NS)
        if clr is not None:
            grad_info = {'type': 'gradient', 'start': clr.get('val', '').upper()}

    # 形状类型
    prstGeom = shape.find('.//a:prstGeom', NS)
    geom_type = prstGeom.get('prst', 'none') if prstGeom is not None else 'none'

    # 文字内容
    texts = shape.findall('.//a:t', NS)
    text_content = [t.text.strip() for t in texts if t.text and t.text.strip()]

    # 段落样式（字号/颜色/粗体）
    runs_info = []
    for rPr in shape.findall('.//a:rPr', NS):
        sz = rPr.get('sz')
        b  = rPr.get('b')
        i  = rPr.get('i')
        # 文字颜色
        tc = rPr.find('a:solidFill/a:srgbClr', NS)
        tc_val = tc.get('val', '').upper() if tc is not None else ''

        run = {
            'size_pt': int(sz) / 100 if sz else None,
            'bold':    b == '1' if b else None,
            'italic':  i == '1' if i else None,
            'color':   tc_val,
        }
        runs_info.append(run)

    # 取最常见的字号/颜色作为段落级别
    sizes  = [r['size_pt'] for r in runs_info if r['size_pt']]
    colors = [r['color']   for r in runs_info if r['color']]
    bolds  = [r['bold']    for r in runs_info if r['bold'] is not None]

    dominant_size  = max(set(sizes),  key=sizes.count)  if sizes  else None
    dominant_color = max(set(colors), key=colors.count) if colors else ''
    has_bold       = any(bolds) if bolds else False

    # 形状种类判断
    if text_content:
        kind = 'text'
    elif geom_type == 'rect' and fill_color:
        kind = 'rect'
    elif geom_type == 'rect':
        kind = 'rect_nofill'
    elif geom_type == 'ellipse':
        kind = 'ellipse'
    elif pic_shapes.__class__.__name__ == 'pic':
        kind = 'image'
    else:
        kind = 'other'

    return {
        'kind':       kind,
        'x_emu':      x_emu,
        'y_emu':      y_emu,
        'w_emu':      cx_emu,
        'h_emu':      cy_emu,
        'x_in':       x_in,
        'y_in':       y_in,
        'w_in':       w_in,
        'h_in':       h_in,
        'fill':       fill_color,
        'fill_grad':  grad_info,
        'geom':       geom_type,
        'texts':      text_content,
        'font_size':  dominant_size,
        'font_color': dominant_color,
        'bold':       has_bold,
        'runs':       runs_info,
    }


# ─────────────────────────────────────────────────────────────────
# 布局分析
# ─────────────────────────────────────────────────────────────────

def analyze_page_layout(page_data: Dict) -> Dict:
    """
    分析单页内容区的布局类型。
    返回 {
        'layout_type': 'card_grid'|'timeline'|'data_table'|
                       'two_column'|'text_section'|'full_media'|'unknown',
        'confidence': 0.0~1.0,
        'reasoning': str,
        'bounds': {x, y, w, h},
        'elements': [element summaries]
    }
    """
    elements = page_data.get('elements', [])
    if not elements:
        return {'layout_type': 'unknown', 'confidence': 0.0,
                'reasoning': 'no elements', 'bounds': {}, 'elements': []}

    # 过滤内容区元素（y > 0.15", y < 7.3", 排除 tiny shapes）
    content = [e for e in elements
               if e['y_in'] > 0.15 and e['y_in'] < 7.3
               and e['w_in'] > 0.3 and e['h_in'] > 0.05]

    rects = [e for e in content if e['kind'] in ('rect', 'rect_nofill') and e['fill']]
    texts = [e for e in content if e['kind'] == 'text']

    # 分析区
    result = {
        'layout_type': 'unknown',
        'confidence':   0.0,
        'reasoning':   '',
        'bounds':      {},
        'elements':     [],
    }

    if not rects and not texts:
        result['reasoning'] = '内容区无有效元素'
        return result

    # 1. 卡片网格检测：2~4个相同尺寸矩形，排列整齐
    grid = detect_card_grid(rects)
    if grid:
        result.update(grid)
        return result

    # 2. 时间线检测：多个矩形左对齐 + y 轴递增
    timeline = detect_timeline(rects, texts)
    if timeline:
        result.update(timeline)
        return result

    # 3. 表格检测：多行等高矩形（row pattern）
    table = detect_table(rects)
    if table:
        result.update(table)
        return result

    # 4. 双栏检测：左右两个大矩形
    two_col = detect_two_column(rects)
    if two_col:
        result.update(two_col)
        return result

    # 5. 文本段落：大量文字元素，矩形少
    if len(texts) >= 3 and len(rects) <= 2:
        result['layout_type'] = 'text_section'
        result['confidence']  = 0.7
        result['reasoning']   = f'文本为主（{len(texts)}个文字元素，{len(rects)}个矩形）'
        return result

    result['reasoning'] = f'未分类: {len(rects)}rects, {len(texts)}texts'
    return result


def detect_card_grid(rects):
    """检测 2×2 / 1×2 / 3×N 卡片网格"""
    if len(rects) < 2:
        return None

    # 按面积分组（相似尺寸为一组）
    areas = [(r['w_in'] * r['h_in'], r) for r in rects]
    areas.sort(key=lambda x: -x[0])  # 大到小

    # 取最大面积的形状作为基准
    base_area = areas[0][0]
    base = areas[0][1]

    # 找同组（面积相差 <20%）
    similar = [r for a, r in areas if abs(a - base_area) / base_area < 0.20]

    if len(similar) < 2:
        return None

    # 行列推断
    xs = sorted(set(round(r['x_in'], 2) for r in similar))
    ys = sorted(set(round(r['y_in'], 2) for r in similar))

    cols = len(xs)
    rows = len(ys)

    if cols >= 2 and rows >= 1:
        # 计算平均尺寸
        avg_w = sum(r['w_in'] for r in similar) / len(similar)
        avg_h = sum(r['h_in'] for r in similar) / len(similar)
        fills = list(set(r['fill'] for r in similar))

        return {
            'layout_type': 'card_grid',
            'confidence':  min(0.9, 0.5 + 0.1 * len(similar)),
            'reasoning':  f'{cols}列×{rows}行卡片网格, {len(similar)}张卡片, 填充色:{fills}',
            'bounds': {
                'x': xs[0], 'y': ys[0],
                'w': round(xs[-1] + avg_w - xs[0], 3),
                'h': round(ys[-1] + avg_h - ys[0], 3),
            },
            'grid': {'cols': cols, 'rows': rows, 'card_count': len(similar),
                     'card_w': round(avg_w, 3), 'card_h': round(avg_h, 3),
                     'xs': xs, 'ys': ys, 'fills': fills},
            'elements': [summarize_elem(r) for r in similar],
        }
    return None


def detect_timeline(rects, texts):
    """检测时间线布局：节点沿 y 轴等间距排列"""
    # 找最小的矩形（dot marker）
    small = [r for r in rects
             if 0.08 <= r['w_in'] <= 0.40
             and 0.08 <= r['h_in'] <= 0.40
             and r['fill']]

    if len(small) < 2:
        return None

    # 按 y 排序
    small.sort(key=lambda r: r['y_in'])
    ys = [r['y_in'] for r in small]

    # 检查 y 间距是否大致相等（允许±30%波动）
    gaps = [ys[i+1] - ys[i] for i in range(len(ys)-1)]
    if not gaps:
        return None

    avg_gap = sum(gaps) / len(gaps)
    uniform = all(abs(g - avg_gap) / avg_gap < 0.30 for g in gaps)

    if uniform and len(small) >= 2:
        return {
            'layout_type': 'timeline',
            'confidence':  min(0.9, 0.6 + 0.1 * len(small)),
            'reasoning':  f'时间线: {len(small)}个节点, 间距≈{avg_gap:.2f}"',
            'bounds': {
                'x': small[0]['x_in'],
                'y': small[0]['y_in'],
                'w': max(r['x_in'] + r['w_in'] for r in small) - small[0]['x_in'],
                'h': small[-1]['y_in'] + small[-1]['h_in'] - small[0]['y_in'],
            },
            'timeline': {
                'nodes': len(small),
                'gap':   round(avg_gap, 3),
                'ys':    [round(y, 3) for y in ys],
                'fills': list(set(r['fill'] for r in small)),
            },
            'elements': [summarize_elem(r) for r in small],
        }
    return None


def detect_table(rects):
    """检测表格布局：多行等高矩形，x 坐标完全对齐"""
    if len(rects) < 3:
        return None

    # 按 y 分组（row）
    rows_data = {}
    for r in rects:
        key_y = round(r['y_in'], 2)
        if key_y not in rows_data:
            rows_data[key_y] = []
        rows_data[key_y].append(r)

    if len(rows_data) < 2:
        return None

    # 每行内的矩形 x 坐标相同
    all_rows_aligned = True
    row_patterns = []

    for y, row_rects in sorted(rows_data.items()):
        if len(row_rects) < 2:
            all_rows_aligned = False
            break
        xs = sorted(round(r['x_in'], 3) for r in row_rects)
        ws = [round(r['w_in'], 3) for r in row_rects]
        row_patterns.append({'xs': xs, 'ws': ws})

    if not all_rows_aligned or not row_patterns:
        return None

    # 检查各行列数是否一致
    first_cols = len(row_patterns[0]['xs'])
    uniform_cols = all(len(rp['xs']) == first_cols for rp in row_patterns)

    if uniform_cols and len(row_patterns) >= 2:
        return {
            'layout_type': 'data_table',
            'confidence':  0.85,
            'reasoning':  f'表格: {first_cols}列×{len(row_patterns)}行',
            'bounds':     {},
            'table': {
                'cols':     first_cols,
                'rows':     len(row_patterns),
                'pattern':  row_patterns,
            },
            'elements': [summarize_elem(r) for r in rects],
        }
    return None


def detect_two_column(rects):
    """检测双栏布局：左右两个大矩形"""
    if len(rects) < 2:
        return None

    # 找最大的两个矩形
    areas = sorted([(r['w_in'] * r['h_in'], r) for r in rects],
                   key=lambda x: -x[0])

    if len(areas) < 2:
        return None

    r1, r2 = areas[0][1], areas[1][1]
    x1, x2 = r1['x_in'], r2['x_in']

    # 判断左右排列：两个矩形 x 不重叠，y 区间重叠
    if abs(x1 - x2) > 1.0:  # x 差距大 = 左右排
        return {
            'layout_type': 'two_column',
            'confidence':   0.8,
            'reasoning':   f'左右双栏: [{x1:.2f}", {x2:.2f}"]',
            'bounds': {
                'x': min(x1, x2),
                'y': min(r1['y_in'], r2['y_in']),
                'w': max(x1 + r1['w_in'], x2 + r2['w_in']) - min(x1, x2),
                'h': max(r1['y_in'] + r1['h_in'], r2['y_in'] + r2['h_in'])
                     - min(r1['y_in'], r2['y_in']),
            },
            'two_column': {
                'left':  {'x': round(x1,3), 'y': round(r1['y_in'],3),
                          'w': round(r1['w_in'],3), 'h': round(r1['h_in'],3)},
                'right': {'x': round(x2,3), 'y': round(r2['y_in'],3),
                          'w': round(r2['w_in'],3), 'h': round(r2['h_in'],3)},
            },
            'elements': [summarize_elem(r) for r in [r1, r2]],
        }
    return None


def summarize_elem(e: Dict) -> Dict:
    """提取元素摘要（用于结构输出）"""
    return {
        'kind':    e['kind'],
        'x':       e['x_in'],
        'y':       e['y_in'],
        'w':       e['w_in'],
        'h':       e['h_in'],
        'fill':    e.get('fill', ''),
        'geom':    e.get('geom', ''),
        'texts':   e.get('texts', []),
        'fs':      e.get('font_size'),
        'color':   e.get('font_color', ''),
        'bold':    e.get('bold'),
    }


# ─────────────────────────────────────────────────────────────────
# 品牌区提取（header / footer）
# ─────────────────────────────────────────────────────────────────

BRAND_AREA_TOP_THRESHOLD    = 0.15  # y < 0.15  → header 区
BRAND_AREA_BOTTOM_THRESHOLD = 7.30  # y > 7.30  → footer 区

def extract_brand_layout(elements: List[Dict]) -> Dict:
    """提取品牌锁定区：顶部条、底部条、标题、页码"""
    brand_rects = [e for e in elements
                  if e['y_in'] < BRAND_AREA_TOP_THRESHOLD
                  and e['kind'] in ('rect', 'rect_nofill')
                  and e['fill']]

    footer_rects = [e for e in elements
                    if e['y_in'] > BRAND_AREA_BOTTOM_THRESHOLD
                    and e['kind'] in ('rect', 'rect_nofill')
                    and e['fill']]

    title_texts = [e for e in elements
                   if 0.15 < e['y_in'] < 1.0
                   and e['kind'] == 'text'
                   and e.get('font_size', 0) >= 24]

    page_nums = [e for e in elements
                 if e['x_in'] > 11.5
                 and e['y_in'] < 1.0
                 and e['kind'] == 'text'
                 and e.get('font_size', 0) and e.get('font_size', 0) <= 14]

    return {
        'header_rects': [summarize_elem(r) for r in brand_rects],
        'footer_rects': [summarize_elem(r) for r in footer_rects],
        'title':        [summarize_elem(t) for t in title_texts],
        'page_number':  [summarize_elem(p) for p in page_nums],
    }


# ─────────────────────────────────────────────────────────────────
# 主题色提取（从实际 shape）
# ─────────────────────────────────────────────────────────────────

def infer_theme_colors(elements: List[Dict]) -> Dict[str, str]:
    """从页面元素反推主题色（最常见的填充色和文字色）"""
    fills  = [e['fill'] for e in elements if e['fill']]
    colors = [e['font_color'] for e in elements if e['font_color']]

    def most_common(lst):
        if not lst: return ''
        return max(set(lst), key=lst.count)

    return {
        'primary_fill':   most_common(fills),
        'primary_text':   most_common(colors),
    }


# ─────────────────────────────────────────────────────────────────
# 主分析流程
# ─────────────────────────────────────────────────────────────────

def analyze_pptx(pptx_path: str) -> Dict:
    """完整分析：raw_data → layout_index"""
    print(f'正在分析: {pptx_path}')

    raw = parse_pptx(pptx_path)
    n   = raw['metadata']['n_slides']

    pages_out = []
    content_layouts = []

    for i, slide in enumerate(raw['slides'], 1):
        print(f'  分析第 {i}/{n} 页...')
        elements = slide.get('elements', [])
        brand    = extract_brand_layout(elements)
        layout   = analyze_page_layout(slide)
        theme    = infer_theme_colors(elements)

        page_out = {
            'index':       i,
            'brand_layout': brand,
            'layout':       layout,
            'theme':        theme,
            'elements':     elements,
        }
        pages_out.append(page_out)
        content_layouts.append({
            'page':   i,
            'type':   layout['layout_type'],
            'conf':   layout['confidence'],
            'reason': layout['reasoning'],
        })

    # 组件建议
    hints = suggest_components(content_layouts)

    return {
        'metadata':         raw['metadata'],
        'pages':           pages_out,
        'content_layouts': content_layouts,
        'component_hints': hints,
    }


def suggest_components(layouts: List[Dict]) -> List[Dict]:
    """根据布局分析给出页型函数建议"""
    hint_map = {
        'card_grid':    'build_card_grid_page',
        'timeline':     'build_timeline_page',
        'data_table':  'build_data_table_page',
        'two_column':   'build_two_column_page',
        'text_section': 'build_text_section_page',
    }
    return [
        {
            'page':  l['page'],
            'type':  l['type'],
            'func':  hint_map.get(l['type'], 'build_card_grid_page'),
            'conf':  l['conf'],
        }
        for l in layouts
        if l['type'] != 'unknown'
    ]


# ─────────────────────────────────────────────────────────────────
# 输出
# ─────────────────────────────────────────────────────────────────

def print_summary(result: Dict):
    """打印摘要"""
    meta = result['metadata']
    print()
    print(f"{'='*60}")
    print(f"参考稿分析完成")
    print(f"{'='*60}")
    print(f"  画布: {meta['canvas_w_in']}\" × {meta['canvas_h_in']}\"")
    print(f"  页数: {meta['n_slides']}")
    print(f"  主题色: {meta['theme_colors']}")
    print()
    print(f"  {'页':>3} | {'布局类型':^18} | {'置信度':^6} | {'原因'}")
    print(f"  {'-'*3}-+-{'-'*20}-+-{'-'*8}- {-20}")
    for l in result['content_layouts']:
        conf = f"{l['conf']:.0%}"
        reason = str(l.get('reason', l.get('reasoning', '')))[:35]
        print(f"  {l['page']:>3} | {l['type']:^18} | {conf:^6} | {reason}")
    print()
    print(f"组件建议:")
    for h in result['component_hints']:
        print(f"  Page {h['page']:>2}: {h['func']}(...)  [{h['type']}]")
    print()


# ─────────────────────────────────────────────────────────────────
# CLI
# ─────────────────────────────────────────────────────────────────

def main():
    ap = argparse.ArgumentParser(description='参考 PPTX → 结构化布局 JSON')
    ap.add_argument('pptx', help='PPTX 文件路径')
    ap.add_argument('-o', '--output', dest='output',
                     help='输出 JSON 路径（默认: <pptx>_layout_index.json）')
    ap.add_argument('--no-hints', action='store_true',
                     help='不生成组件建议')
    args = ap.parse_args()

    result = analyze_pptx(args.pptx)

    if args.no_hints:
        result.pop('component_hints', None)

    out_path = args.output or str(Path(args.pptx).with_suffix('')) + '_layout_index.json'
    with open(out_path, 'w', encoding='utf-8') as f:
        json.dump(result, f, ensure_ascii=False, indent=2)

    print(f'\n✓ 结构化数据已保存: {out_path}')
    print_summary(result)


if __name__ == '__main__':
    main()
