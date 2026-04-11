#!/usr/bin/env python3
"""Import a reference PPT into case_library and generate structured analysis."""

from __future__ import annotations

import argparse
import json
import re
import shutil
from collections import Counter
from pathlib import Path
from typing import Any

from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
CASE_LIBRARY = ROOT / "case_library"


def slugify(value: str) -> str:
    value = re.sub(r"[^0-9A-Za-z\u4e00-\u9fff_-]+", "_", value.strip())
    value = re.sub(r"_+", "_", value).strip("_")
    return value or "case"


def extract_slide_title(slide) -> str:
    texts: list[tuple[float, str]] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = "\n".join(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()).strip()
        if not text:
            continue
        size = 0.0
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    size = max(size, float(run.font.size.pt))
        texts.append((size, text.splitlines()[0][:80]))
    if not texts:
        return "未识别标题"
    texts.sort(key=lambda item: (-item[0], item[1]))
    return texts[0][1]


def analyze_ppt(path: Path) -> dict[str, Any]:
    prs = Presentation(str(path))
    slides: list[dict[str, Any]] = []
    shape_counter: Counter[str] = Counter()

    for idx, slide in enumerate(prs.slides, start=1):
        shape_types: Counter[str] = Counter()
        text_samples: list[str] = []
        text_boxes = 0
        image_count = 0
        table_count = 0
        for shape in slide.shapes:
            shape_name = str(getattr(shape, "shape_type", "unknown"))
            shape_counter[shape_name] += 1
            shape_types[shape_name] += 1
            if getattr(shape, "has_text_frame", False):
                text_boxes += 1
                text = "\n".join(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()).strip()
                if text:
                    text_samples.append(text.splitlines()[0][:100])
            if getattr(shape, "shape_type", None) == 13:
                image_count += 1
            if getattr(shape, "shape_type", None) == 19:
                table_count += 1
        slides.append(
            {
                "slide_num": idx,
                "title_guess": extract_slide_title(slide),
                "shape_count": len(slide.shapes),
                "text_boxes": text_boxes,
                "images": image_count,
                "tables": table_count,
                "top_texts": text_samples[:5],
                "shape_distribution": dict(shape_types.most_common(8)),
            }
        )

    return {
        "file": str(path),
        "slide_count": len(slides),
        "shape_distribution": dict(shape_counter.most_common(12)),
        "slides": slides,
    }


def write_markdown(case_dir: Path, analysis: dict[str, Any], domain: str, case_name: str, source_name: str) -> None:
    slides = analysis["slides"]
    titles = [f"第{s['slide_num']}页：{s['title_guess']}" for s in slides]
    dense = sorted(slides, key=lambda s: (-s["shape_count"], -s["text_boxes"]))[:5]

    (case_dir / "case_meta.md").write_text(
        f"# 案例元信息\n\n- 案例名称：{case_name}\n- 行业：{domain}\n- 来源文件：{source_name}\n- 总页数：{analysis['slide_count']}\n",
        encoding="utf-8",
    )
    (case_dir / "deck_outline.md").write_text(
        "# Deck 结构\n\n" + "\n".join(f"- {line}" for line in titles),
        encoding="utf-8",
    )
    (case_dir / "page_patterns.md").write_text(
        "# 页面模式\n\n"
        + "\n".join(
            f"- 第{s['slide_num']}页《{s['title_guess']}》：形状 {s['shape_count']}，文本框 {s['text_boxes']}，图片 {s['images']}，表格 {s['tables']}"
            for s in slides[:15]
        ),
        encoding="utf-8",
    )
    (case_dir / "diagram_patterns.md").write_text(
        "# 图形模式\n\n"
        + "\n".join(f"- 形状类型 {shape}: {count}" for shape, count in analysis["shape_distribution"].items())
        + "\n\n## 高密度页面\n\n"
        + "\n".join(
            f"- 第{s['slide_num']}页《{s['title_guess']}》：形状 {s['shape_count']}，文本框 {s['text_boxes']}"
            for s in dense
        ),
        encoding="utf-8",
    )
    (case_dir / "writing_logic.md").write_text(
        "# 文本逻辑\n\n"
        + "\n".join(
            f"- 第{s['slide_num']}页《{s['title_guess']}》文本样例：" + " / ".join(s["top_texts"][:3])
            for s in slides[:15]
        ),
        encoding="utf-8",
    )
    (case_dir / "soft_issue_risks.md").write_text(
        "# 潜在软性风险\n\n"
        "- 是否存在标题与正文主线不一致的页面\n"
        "- 是否存在复杂结构但证据悬空的页面\n"
        "- 是否存在页尾没有管理收束的页面\n",
        encoding="utf-8",
    )
    (case_dir / "merge_suggestions.md").write_text(
        "# 合并建议\n\n"
        "- 建议区分哪些是模板骨架规则\n"
        "- 建议区分哪些是行业表达规则\n"
        "- 建议区分哪些仅是单案例技巧\n",
        encoding="utf-8",
    )


def update_index(domain: str, case_name: str, case_dir: Path) -> None:
    index_path = CASE_LIBRARY / "index.json"
    data = json.loads(index_path.read_text(encoding="utf-8"))
    domains = data.setdefault("domains", {})
    cases = domains.setdefault(domain, [])
    record = {"name": case_name, "path": str(case_dir.relative_to(ROOT))}
    if record not in cases:
        cases.append(record)
    index_path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser(description="Ingest a reference PPT into case_library.")
    parser.add_argument("pptx", help="Path to source PPTX file")
    parser.add_argument("--domain", default="general", help="Domain bucket, e.g. security_service")
    parser.add_argument("--case-name", help="Override case folder name")
    parser.add_argument("--copy-source", action="store_true", help="Copy source file into source_files/")
    args = parser.parse_args()

    pptx_path = Path(args.pptx).expanduser().resolve()
    if not pptx_path.exists():
        raise SystemExit(f"PPTX not found: {pptx_path}")

    domain = slugify(args.domain)
    case_name = slugify(args.case_name or pptx_path.stem)
    case_dir = CASE_LIBRARY / domain / case_name
    case_dir.mkdir(parents=True, exist_ok=True)
    (case_dir / "source_files").mkdir(exist_ok=True)

    if args.copy_source:
        shutil.copy2(pptx_path, case_dir / "source_files" / pptx_path.name)

    analysis = analyze_ppt(pptx_path)
    (case_dir / "analysis.json").write_text(json.dumps(analysis, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    write_markdown(case_dir, analysis, domain, case_name, pptx_path.name)
    update_index(domain, case_name, case_dir)

    print(f"Case ingested: {case_dir}")


if __name__ == "__main__":
    main()
